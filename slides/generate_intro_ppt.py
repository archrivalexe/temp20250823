import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --------------------------
# 配置
# --------------------------
DEFAULT_FONT_NAME = "Microsoft YaHei"
TITLE_FONT_SIZE = Pt(40)
SUBTITLE_FONT_SIZE = Pt(28)
BODY_FONT_SIZE = Pt(26)
CAPTION_FONT_SIZE = Pt(18)

COLORS = {
    "blue": RGBColor(37, 99, 235),   # #2563EB 教师
    "green": RGBColor(22, 163, 74),  # #16A34A 学生
    "gray": RGBColor(107, 114, 128), # #6B7280 观察者/说明
    "black": RGBColor(17, 24, 39)    # #111827 文本
}

OUTPUT_DIR = "slides"
OUTPUT_FILENAME = os.path.join(OUTPUT_DIR, "AIEd_Student_Social_Cognition_Intro_CN.pptx")

# --------------------------
# 文案（可按需修改）
# --------------------------

TITLE_TEXT = "学生对教育者AI使用行为的社会认知\n自我报告与源记忆的证据"
SUBTITLE_TEXT = "当老师用AI时，学生在想什么？"

SLIDE2_TITLE = "背景：AIEd 的普及与关注点"
SLIDE2_BULLETS = [
    "AIEd 在课堂与评价中加速普及",
    "研究多关注教学促进与效率提升",
    "自适应学习、即时反馈、负担减轻等"
]
NOTES_2 = (
    "用 1–2 个熟悉的 AIEd 场景作例子（如作业反馈、备课生成）。"
    "先给出积极语气，为后文的阻碍与挑战做衔接。"
)

SLIDE3_TITLE = "教育中的阻碍更大"
SLIDE3_BULLETS = [
    "相比其他领域，教育阻碍更明显（Bates et al., 2020）",
    "核心挑战：学生的接受与适应（Renz & Hilbig, 2020）",
    "不仅是技术，更是社会心理与信任问题"
]
NOTES_3 = (
    "强调“教育不是纯技术问题”，涉及规范、伦理、信任与身份角色。"
    "此页仅提作者与年份，不展开方法细节。"
)

SLIDE4_TITLE = "研究空白：人-人关系如何被AI改变？"
SLIDE4_BULLETS = [
    "大量研究：关注人-机关系",
    "较少研究：AI使用如何改变人-人关系",
    "课堂里，AI介入学生-教师互动成为新变量"
]
NOTES_4 = "点出社会知觉的入口：学生如何看待使用AI的教师（能力、公平、可信度、关怀等）。"

SLIDE5_TITLE = "角色生态：三类角色，一张图"
SLIDE5_CAPTION = "学生（被动） → 对教师（主动）的社会知觉；以往多为观察者视角（组织情境）"
NOTES_5 = (
    "图示颜色建议：教师（蓝）、学生（绿）、观察者（灰）。"
    "箭头高亮本研究“学生→教师”的知觉形成路径。"
)

SLIDE6_TITLE = "现有研究的视角与局限"
SLIDE6_BULLETS = [
    "现有少量研究多从观察者视角（组织情境）",
    "关注同事/上级对员工能力的感知",
    "教育情境中，学生才是被动使用者"
]
NOTES_6 = "对比：组织（观察者→使用者） vs 教育（被动使用者→主动使用者）。"

SLIDE7_TITLE = "本研究切入：学生的社会知觉"
SLIDE7_BULLETS = [
    "目标：从学生视角理解教师AI使用的社会影响",
    "关注：对教师的能力、公平性、可信度等知觉",
    "方法：自我报告与源记忆证据（互补）"
]
NOTES_7 = "自我报告＝主观评价；源记忆＝对信息来源辨识的客观基础。"

SLIDE8_TITLE = "关键研究问题"
SLIDE8_BULLETS = [
    "RQ1 学生如何基于教师的AI使用（透明度/频率/情境）形成对教师能力与公平的知觉？",
    "RQ2 学生的源记忆准确性与其社会知觉是否相关？",
    "RQ3 不同教学情境（讲授/作业反馈/评估）是否存在模式差异？"
]
NOTES_8 = "篇幅有限可保留 RQ1、RQ2；变量操作化细节放方法部分。"

SLIDE9_TITLE = "概念澄清"
SLIDE9_BULLETS = [
    "社会认知/社会知觉：对他人的特质与意图的推断",
    "自我报告：学生主观评价与感受",
    "源记忆：对信息来源的记忆与辨识"
]
NOTES_9 = "术语门槛一页搞定，降低后续理解成本。每条保持 1 行。"

SLIDE10_TITLE = "贡献与意义"
SLIDE10_BULLETS = [
    "视角创新：从观察者转向被动使用者（学生）",
    "机制探索：连接源记忆与社会知觉",
    "实践启示：为教师AI使用的透明度与沟通策略提供依据"
]
NOTES_10 = "用 1 句话总结贡献与预期影响，便于记忆。"

SLIDE_REF_TITLE = "参考文献（展示友好格式）"
SLIDE_REF_BULLETS = [
    "Bates et al. (2020). [请核对完整题目与期刊]",
    "Renz & Hilbig (2020). [请核对完整题目与期刊]"
]
NOTES_REF = "PPT中仅列作者与年份；完整参考文献可放在备份页或论文。"

# --------------------------
# 工具函数
# --------------------------

def set_text_format(text_frame, font_size=BODY_FONT_SIZE, bold=False, align=PP_ALIGN.LEFT):
    text_frame.word_wrap = True
    text_frame.auto_size = False
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    for p in text_frame.paragraphs:
        for r in p.runs:
            r.font.name = DEFAULT_FONT_NAME
            r.font.size = font_size
            r.font.bold = bold
            r.font.color.rgb = COLORS["black"]
        p.alignment = align


def add_speaker_notes(slide, notes_text: str):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    tf.text = notes_text


def add_title_slide(prs, title_text, subtitle_text):
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    set_text_format(title.text_frame, font_size=TITLE_FONT_SIZE, bold=True, align=PP_ALIGN.LEFT)

    subtitle.text = subtitle_text
    set_text_format(subtitle.text_frame, font_size=SUBTITLE_FONT_SIZE, bold=False, align=PP_ALIGN.LEFT)
    return slide


def add_bullet_slide(prs, title_text, bullets, notes_text=""):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    set_text_format(slide.shapes.title.text_frame, font_size=TITLE_FONT_SIZE, bold=True)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = b
        for r in p.runs:
            r.font.name = DEFAULT_FONT_NAME
            r.font.size = BODY_FONT_SIZE
            r.font.color.rgb = COLORS["black"]
        p.level = 0

    if notes_text:
        add_speaker_notes(slide, notes_text)
    return slide


def add_role_ecology_slide(prs, title_text, caption_text, notes_text=""):
    # Blank layout for custom drawing
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(1.0))
    tf = title_box.text_frame
    tf.text = title_text
    set_text_format(tf, font_size=TITLE_FONT_SIZE, bold=True)

    # Positions
    cx, cy = Inches(2.6), Inches(3.8)   # Student (left)
    tx, ty = Inches(9.2), Inches(3.8)   # Teacher (right)
    ox, oy = Inches(5.9), Inches(1.6)   # Observer (top)

    radius = Inches(1.2)

    # Student node
    student = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - radius/2, cy - radius/2, radius, radius)
    student.fill.solid()
    student.fill.fore_color.rgb = COLORS["green"]
    student.line.color.rgb = COLORS["green"]
    s_tf = student.text_frame
    s_tf.text = "学生\n（被动使用者）"
    set_text_format(s_tf, font_size=Pt(20), bold=True, align=PP_ALIGN.CENTER)

    # Teacher node
    teacher = slide.shapes.add_shape(MSO_SHAPE.OVAL, tx - radius/2, ty - radius/2, radius, radius)
    teacher.fill.solid()
    teacher.fill.fore_color.rgb = COLORS["blue"]
    teacher.line.color.rgb = COLORS["blue"]
    t_tf = teacher.text_frame
    t_tf.text = "教师\n（AI主动使用者）"
    set_text_format(t_tf, font_size=Pt(20), bold=True, align=PP_ALIGN.CENTER)

    # Observer node
    observer = slide.shapes.add_shape(MSO_SHAPE.OVAL, ox - radius/2, oy - radius/2, radius, radius)
    observer.fill.solid()
    observer.fill.fore_color.rgb = COLORS["gray"]
    observer.line.color.rgb = COLORS["gray"]
    o_tf = observer.text_frame
    o_tf.text = "观察者"
    set_text_format(o_tf, font_size=Pt(20), bold=True, align=PP_ALIGN.CENTER)

    # Arrow: Student -> Teacher
    arrow = slide.shapes.add_connector(1, cx + radius/2, cy, tx - radius/2, ty)  # straight connector
    arrow.line.color.rgb = COLORS["blue"]
    arrow.line.width = Inches(0.08)

    # Label for main path
    label_box = slide.shapes.add_textbox(Inches((2.6+9.2)/2 - 1.2), Inches(cy - 0.9), Inches(2.4), Inches(0.6))
    ltf = label_box.text_frame
    ltf.text = "社会知觉"
    set_text_format(ltf, font_size=Pt(18), bold=True, align=PP_ALIGN.CENTER)

    # Arrow: Observer -> Teacher (prior literature)
    obs_to_teacher = slide.shapes.add_connector(1, ox, oy + radius/2, tx, ty - radius/2)
    obs_to_teacher.line.color.rgb = COLORS["gray"]
    obs_to_teacher.line.width = Inches(0.04)

    # AI tools box near Teacher
    ai_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx + Inches(0.9), ty - Inches(0.5), Inches(2.0), Inches(1.0))
    ai_box.fill.solid()
    ai_box.fill.fore_color.rgb = RGBColor(229, 231, 235) # light gray background
    ai_box.line.color.rgb = COLORS["gray"]
    ai_tf = ai_box.text_frame
    ai_tf.text = "AI 工具与流程"
    set_text_format(ai_tf, font_size=Pt(16), bold=False, align=PP_ALIGN.CENTER)

    # Arrow from AI box to Teacher
    ai_to_teacher = slide.shapes.add_connector(1, tx + Inches(0.9), ty, tx + radius/2, ty)
    ai_to_teacher.line.color.rgb = COLORS["blue"]
    ai_to_teacher.line.width = Inches(0.04)

    # Caption
    cap_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(12), Inches(0.8))
    cap_tf = cap_box.text_frame
    cap_tf.text = caption_text
    set_text_format(cap_tf, font_size=CAPTION_FONT_SIZE, bold=False, align=PP_ALIGN.LEFT)

    if notes_text:
        add_speaker_notes(slide, notes_text)
    return slide


def add_references_slide(prs, title_text, bullets, notes_text=""):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    set_text_format(slide.shapes.title.text_frame, font_size=TITLE_FONT_SIZE, bold=True)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        for r in p.runs:
            r.font.name = DEFAULT_FONT_NAME
            r.font.size = Pt(22)
            r.font.color.rgb = COLORS["black"]
        p.level = 0

    if notes_text:
        add_speaker_notes(slide, notes_text)
    return slide

# --------------------------
# 构建演示文稿
# --------------------------

def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    prs = Presentation()
    # 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    s1 = add_title_slide(prs, TITLE_TEXT, SUBTITLE_TEXT)

    # Slide 2
    s2 = add_bullet_slide(prs, SLIDE2_TITLE, SLIDE2_BULLETS, NOTES_2)

    # Slide 3
    s3 = add_bullet_slide(prs, SLIDE3_TITLE, SLIDE3_BULLETS, NOTES_3)

    # Slide 4
    s4 = add_bullet_slide(prs, SLIDE4_TITLE, SLIDE4_BULLETS, NOTES_4)

    # Slide 5 (role ecology diagram)
    s5 = add_role_ecology_slide(prs, SLIDE5_TITLE, SLIDE5_CAPTION, NOTES_5)

    # Slide 6
    s6 = add_bullet_slide(prs, SLIDE6_TITLE, SLIDE6_BULLETS, NOTES_6)

    # Slide 7
    s7 = add_bullet_slide(prs, SLIDE7_TITLE, SLIDE7_BULLETS, NOTES_7)

    # Slide 8
    s8 = add_bullet_slide(prs, SLIDE8_TITLE, SLIDE8_BULLETS, NOTES_8)

    # Slide 9
    s9 = add_bullet_slide(prs, SLIDE9_TITLE, SLIDE9_BULLETS, NOTES_9)

    # Slide 10
    s10 = add_bullet_slide(prs, SLIDE10_TITLE, SLIDE10_BULLETS, NOTES_10)

    # Slide 11 (References)
    s11 = add_references_slide(prs, SLIDE_REF_TITLE, SLIDE_REF_BULLETS, NOTES_REF)

    prs.save(OUTPUT_FILENAME)
    print(f"Saved: {OUTPUT_FILENAME}")


if __name__ == "__main__":
    main()
