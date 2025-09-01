#!/usr/bin/env python3
# -*- coding: utf-8 -*-

# Fix for Python 3.10+ compatibility with older python-pptx
import collections
import collections.abc
if not hasattr(collections, 'Container'):
    collections.Container = collections.abc.Container
if not hasattr(collections, 'Mapping'):
    collections.Mapping = collections.abc.Mapping
if not hasattr(collections, 'MutableMapping'):
    collections.MutableMapping = collections.abc.MutableMapping
if not hasattr(collections, 'Sequence'):
    collections.Sequence = collections.abc.Sequence

import argparse
import os
import re
from collections import defaultdict, OrderedDict

from pptx import Presentation
from docx import Document
from docx.shared import Pt

US_KEYS = [
    "美国", "U.S.", "US", "USA", "U.S.A", "United States", "America"
]
UK_KEYS = [
    "英国", "UK", "U.K.", "United Kingdom", "Great Britain", "Britain",
    "England", "Scotland", "Wales", "Northern Ireland", "GB"
]

FIELD_PATTERNS = OrderedDict([
    ("name", [r"^\s*(项目名称|Program Name|Programme|Project|项目)[:：]\s*(.*)$"]),
    ("duration", [r"^\s*(学制|Duration|Length|Program Length|Length of Study)[:：]?\s*(.*)$"]),
    ("fee", [r"^\s*(学费|费用|Tuition|Tuition Fee|Fee|Tuition Fees)[:：]?\s*(.*)$"]),
    ("intro", [r"^\s*(项目介绍|简介|Overview|Program Overview|Description|About)[:：]?\s*(.*)$"]),
])

ENTRY_START_HINTS = [
    r"^\s*[•·\-—–]\s+",
    r"^\s*(项目名称|Program Name|Programme|项目)[:：]\s+",
]

def detect_country(text: str) -> str:
    t = (text or "").lower()
    for k in US_KEYS:
        if k.lower() in t:
            return "美国"
    for k in UK_KEYS:
        if k.lower() in t:
            return "英国"
    return "其他国家"

def extract_slide_texts(pptx_path: str):
    prs = Presentation(pptx_path)
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title and hasattr(slide.shapes.title, "text"):
            title = (slide.shapes.title.text or "").strip()
        lines = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    text = (para.text or "").strip()
                if text:
                    lines.append(text)
        deduped = []
        for l in lines:
            if not deduped or deduped[-1] != l:
                deduped.append(l)
        slides.append({
            "index": idx,
            "title": title,
            "lines": deduped
        })
    return slides

def split_entries(lines):
    entries = []
    cur = []
    start_regex = re.compile("|".join(ENTRY_START_HINTS))
    for line in lines:
        if start_regex.search(line) and cur:
            entries.append(cur)
            cur = [line]
        else:
            cur.append(line)
    if cur:
        entries.append(cur)
    if len(entries) >= 6 and sum(len(e) for e in entries) / len(entries) < 3:
        return [lines]
    return entries

def parse_entry(lines):
    entry = {"name": "", "duration": "", "fee": "", "intro": "", "raw": []}
    entry["raw"] = lines[:]
    remaining = []
    for line in lines:
        matched_any = False
        for field, patterns in FIELD_PATTERNS.items():
            for pat in patterns:
                m = re.match(pat, line, flags=re.IGNORECASE)
                if m:
                    val = m.group(2).strip() if m.lastindex and m.lastindex >= 2 else ""
                    if not entry[field]:
                        entry[field] = val
                        matched_any = True
                    break
            if matched_any:
                break
        if not matched_any:
            remaining.append(line)

    if not entry["name"]:
        entry["name"] = next((l for l in remaining if l.strip()), "")

    if not entry["intro"]:
        intro_candidates = []
        for l in remaining:
            if re.search(r"(学费|费用|Tuition|Fee)", l, flags=re.IGNORECASE):
                continue
            if re.search(r"(学制|Duration|Length)", l, flags=re.IGNORECASE):
                continue
            intro_candidates.append(l)
        if intro_candidates and entry["name"]:
            if intro_candidates[0].strip() == entry["name"].strip():
                intro_candidates = intro_candidates[1:]
        entry["intro"] = "\n".join(intro_candidates[:6]).strip()

    return entry

def ensure_parent_dir(path: str):
    d = os.path.dirname(path)
    if d and not os.path.exists(d):
        os.makedirs(d, exist_ok=True)


def build_doc(slide_groups, output_path: str):
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "微软雅黑"
    style.font.size = Pt(11)

    doc.add_heading("国内外著名大学心理学系硕士项目费用分析", level=0)
    doc.add_paragraph("说明：以下内容来自提供的 PPT 文档自动抽取与整理，仅展示学费等原币种费用，未包含生活费/杂费。按国家分组为'美国 / 英国 / 其他国家'，并尽量保持原始顺序。")

    for country in ["美国", "英国", "其他国家"]:
        group = slide_groups.get(country, [])
        if not group:
            continue
        doc.add_heading(country, level=1)
        for item in group:
            slide_title = item["slide_title"]
            for entry in item["entries"]:
                name = entry["name"] or slide_title or "未命名项目"
                doc.add_heading(name, level=2)
                if entry["duration"]:
                    p = doc.add_paragraph()
                    p.add_run("学制：").bold = True
                    p.add_run(entry["duration"])
                if entry["fee"]:
                    p = doc.add_paragraph()
                    p.add_run("费用：").bold = True
                    p.add_run(entry["fee"])
                if entry["intro"]:
                    p = doc.add_paragraph()
                    p.add_run("项目介绍：").bold = True
                    p.add_run("\n" + entry["intro"])
                if not (entry["duration"] or entry["fee"] or entry["intro"]):
                    p = doc.add_paragraph()
                    p.add_run("项目信息（自动抽取，原始文本）：").bold = True
                    p.add_run("\n" + "\n".join(entry["raw"]))

    ensure_parent_dir(output_path)
    doc.save(output_path)


def main():
    parser = argparse.ArgumentParser(description="从 PPTX 抽取心理学硕士项目信息并生成 DOCX（原币种、按国家分组）")
    parser.add_argument("--input", "-i", required=True, help="输入的 PPTX 文件路径")
    parser.add_argument("--output", "-o", required=True, help="输出的 DOCX 文件路径")
    args = parser.parse_args()

    slides = extract_slide_texts(args.input)

    grouped = defaultdict(list)
    for s in slides:
        country = detect_country(s["title"])
        if country == "其他国家":
            body_text = " ".join(s["lines"])
            country = detect_country(body_text)
        entries_blocks = split_entries(s["lines"])
        entries = [parse_entry(block) for block in entries_blocks]
        grouped[country].append({
            "slide_index": s["index"],
            "slide_title": s["title"],
            "entries": entries
        })

    build_doc(grouped, args.output)

if __name__ == "__main__":
    main()